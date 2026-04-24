import io
import json
import re
from app.ai.mistral_client import chat


def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def extract_text_from_pdf(file_bytes: bytes) -> str:
    import pdfplumber
    lines = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                lines.append(text.strip())
    return "\n".join(lines)


def extract_text_from_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    lines = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    return "\n".join(lines)


def extract_project_fields(raw_text: str) -> dict:
    prompt = f"""You are extracting structured fields from a student project document.

Document text:
\"\"\"
{raw_text[:4000]}
\"\"\"

Extract and return ONLY valid JSON with these fields:
{{
  "title": "project title (short, 5-100 chars)",
  "description": "project description (at least 50 chars, summarize if needed)",
  "modules": ["list", "of", "project", "modules", "or", "features"],
  "technologies": ["list", "of", "technologies", "tools", "languages", "used"],
  "team_members": ["list of team member names if found, else empty array"],
  "domain": "one of: Machine Learning, Web Development, Mobile App, IoT, Cybersecurity, Data Analytics, Blockchain, Cloud Computing, AR/VR, Other"
}}

Rules:
- title must be a clean project name, not a slide heading like "Introduction"
- description must be a coherent paragraph, not bullet points
- modules must have at least 2 items; infer from features/chapters if not explicit
- technologies must have at least 1 item; detect from context (frameworks, languages, DBs)
- Return ONLY the JSON object, no markdown, no explanation"""

    response = chat(prompt, system="You are a JSON extraction engine. Return only valid JSON.")
    match = re.search(r"\{.*\}", response, re.DOTALL)
    if not match:
        raise ValueError("Mistral did not return valid JSON")
    return json.loads(match.group())


def parse_document(filename: str, file_bytes: bytes) -> dict:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext in ("pptx", "ppt"):
        raw_text = extract_text_from_pptx(file_bytes)
    elif ext == "pdf":
        raw_text = extract_text_from_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        raw_text = extract_text_from_docx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if len(raw_text.strip()) < 50:
        raise ValueError("Could not extract enough text from document")

    return extract_project_fields(raw_text)
